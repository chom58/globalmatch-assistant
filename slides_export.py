"""CV提案コメントをGoogleスライド貼付用の.pptxに変換する。

「企業の価値をデザインする。」(Value Create) のデザインシステムに準拠した
ダーク基調のスライドテンプレート:
- 背景: 純黒 (#000000 = surface.base)
- カード: ダークグレー (#313131 = surface.raised) + 淡い境界 (#cccccc = border.muted)
- 本文・見出し: 白 (#ffffff = text.secondary)
- アクセント: バリュークリエイト ブランドイエロー (#F7CA45)
- フォント: Noto Sans Japanese（fallback で PowerPoint 既定和文フォントへ）

レイアウト: 16:9 ワイドスクリーン
- 上部: ダークヘッダー + Headline（イエローアンダーライン）
- 中段: 2x2 グリッド (Career / Strengths / Education / Assessment)
- 右下: Value Create フッター
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# === デザイントークン（「企業の価値をデザインする。」準拠のライトテーマ） ===
# 暗い背景は商談資料で読みづらい指摘があったため、ベースを白に反転。
# イエローアクセント + Noto Sans Japanese はブランド要素として維持。
BRAND_YELLOW = RGBColor(0xF7, 0xCA, 0x45)           # accent (ブランドイエロー)
SURFACE_BASE = RGBColor(0xFF, 0xFF, 0xFF)           # 背景: 純白
SURFACE_RAISED = RGBColor(0xFA, 0xFA, 0xFA)         # カード: ごく薄いグレー
SURFACE_HEADER = RGBColor(0xF5, 0xF5, 0xF5)         # ヘッダー帯: やや濃いめの薄グレー
TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x1A)           # 本文・見出し
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)             # フッター等の補助テキスト
BORDER_MUTED = RGBColor(0xCC, 0xCC, 0xCC)           # color.border.muted

# Font family: Noto Sans Japanese（実行環境に無ければ PowerPoint 既定和文フォントへ fallback）
FONT_FAMILY = "Noto Sans Japanese"

# 16:9 ワイドスクリーン
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# セクションタイトルの日英表示（ヘッダー内部で用いる）
SECTION_LABELS_JA = {
    "Headline": "Headline",
    "Career": "Career / 経歴",
    "Strengths": "Strengths / 強み",
    "Education / Research": "Education / Research",
    "Assessment": "Assessment / 評価",
}
SECTION_LABELS_EN = {
    "Headline": "Headline",
    "Career": "Career",
    "Strengths": "Strengths",
    "Education / Research": "Education / Research",
    "Assessment": "Assessment",
}


@dataclass
class CVSection:
    """抽出済みセクションの正規化データ。"""

    headline: str = ""
    career: str = ""
    strengths: str = ""
    education: str = ""
    assessment: str = ""

    @classmethod
    def from_pairs(cls, pairs: list[tuple[str, str]]) -> "CVSection":
        """_parse_cv_sections の結果 [(name, body), ...] から構築する。
        見出しの数字プレフィックス・大文字小文字・日本語副題を許容する。"""
        result = cls()
        for name, body in pairs:
            key = _normalize_section_name(name)
            if key == "headline":
                result.headline = body
            elif key == "career":
                result.career = body
            elif key == "strengths":
                result.strengths = body
            elif key == "education":
                result.education = body
            elif key == "assessment":
                result.assessment = body
        return result


_SECTION_NORMALIZE = {
    "headline": "headline",
    "career": "career",
    "strengths": "strengths",
    "education": "education",
    "education / research": "education",
    "research": "education",
    "assessment": "assessment",
}


def _normalize_section_name(raw: str) -> str:
    """'1. 👤 Headline' / '1. Headline' / 'Education / Research' などを正規化キーに変換する。
    先頭の数字プレフィックス・絵文字・空白を除去してからキー照合する。"""
    # 先頭の数字とドット・空白を除去
    cleaned = re.sub(r"^\s*\d+\.\s*", "", raw).strip()
    # 先頭の絵文字・記号（非ASCII英数字）を除去
    cleaned = re.sub(r"^[^A-Za-z]+", "", cleaned).strip().lower()
    # キーとの部分一致
    for key, norm in _SECTION_NORMALIZE.items():
        if cleaned.startswith(key):
            return norm
    return cleaned


def build_cv_proposal_pptx(
    candidates: list[list[tuple[str, str]]],
    language: str = "ja",
    labels: list[str] | None = None,
) -> bytes:
    """候補者ごとにスライド1枚を生成して1つの.pptxにまとめて返す。

    Args:
        candidates: 候補者ごとの [(section_name, body), ...] のリスト
        language: "ja" | "en"
        labels: 候補者識別子（スライド左上に小さく表示、任意）

    Returns:
        .pptx のバイナリ
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # 完全な空白レイアウト

    section_labels = SECTION_LABELS_JA if language == "ja" else SECTION_LABELS_EN

    if not candidates:
        candidates = [[]]  # 空でも1枚は出す

    labels = labels or [""] * len(candidates)
    if len(labels) < len(candidates):
        labels = list(labels) + [""] * (len(candidates) - len(labels))

    for idx, (pairs, label) in enumerate(zip(candidates, labels)):
        slide = prs.slides.add_slide(blank_layout)
        data = CVSection.from_pairs(pairs)
        _render_slide(slide, data, language, section_labels, label, idx + 1, len(candidates))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _render_slide(
    slide,
    data: "CVSection",
    language: str,
    section_labels: dict,
    label: str,
    slide_no: int,
    total: int,
) -> None:
    # --- 0. スライド全体の背景（黒） ---
    # slide.background.fill を solid black に設定することで、Keynote/PowerPoint 両方で
    # 確実にダーク基調を適用する（add_shape で敷くと Keynote 描画で無視される環境があるため）
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = SURFACE_BASE

    # --- 1. 上部ヘッダー帯（薄グレー + 下端イエローライン） ---
    header_h = Inches(0.55)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, header_h
    )
    header.line.fill.background()
    header.fill.solid()
    header.fill.fore_color.rgb = SURFACE_HEADER
    title_text = "Candidate Pitch" if language == "en" else "候補者プロフィール"
    _set_shape_text(
        header,
        title_text,
        font_size=18,
        bold=True,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
        left_margin=Inches(0.4),
    )
    # ヘッダー下端のイエロー 2px ライン（ブランド識別）
    header_underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, header_h - Inches(0.03), SLIDE_W, Inches(0.03)
    )
    header_underline.line.fill.background()
    header_underline.fill.solid()
    header_underline.fill.fore_color.rgb = BRAND_YELLOW

    # 右側にページ番号 & label
    right_text_parts = []
    if label:
        right_text_parts.append(label)
    if total > 1:
        right_text_parts.append(f"{slide_no} / {total}")
    if right_text_parts:
        right_tb = slide.shapes.add_textbox(
            SLIDE_W - Inches(4.0), Inches(0.1), Inches(3.6), Inches(0.35)
        )
        _set_textbox_text(
            right_tb,
            "  ".join(right_text_parts),
            font_size=11,
            bold=False,
            color=TEXT_PRIMARY,
            align=PP_ALIGN.RIGHT,
        )

    # --- 2. Headline（ヘッダー直下、全幅） ---
    headline_top = Inches(0.78)
    headline_height = Inches(0.85)
    headline_box = slide.shapes.add_textbox(
        Inches(0.5), headline_top, SLIDE_W - Inches(1.0), headline_height
    )
    _set_textbox_text(
        headline_box,
        data.headline or "—",
        font_size=22,
        bold=True,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Headline 下のイエロー アンダーライン
    ul = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        headline_top + headline_height,
        Inches(2.2),
        Inches(0.05),
    )
    ul.line.fill.background()
    ul.fill.solid()
    ul.fill.fore_color.rgb = BRAND_YELLOW

    # --- 3. 2x2 グリッド ---
    grid_top = Inches(1.9)
    grid_left = Inches(0.5)
    grid_right = SLIDE_W - Inches(0.5)
    grid_bottom = SLIDE_H - Inches(0.55)  # フッター分を除く
    gap = Inches(0.22)

    total_w = grid_right - grid_left
    total_h = grid_bottom - grid_top
    cell_w = (total_w - gap) / 2
    cell_h = (total_h - gap) / 2

    cells = [
        (0, 0, "Career", data.career),
        (0, 1, "Strengths", data.strengths),
        (1, 0, "Education / Research", data.education),
        (1, 1, "Assessment", data.assessment),
    ]

    for row, col, key, body in cells:
        x = grid_left + col * (cell_w + gap)
        y = grid_top + row * (cell_h + gap)
        title = section_labels.get(key, key)
        _draw_cell(slide, x, y, cell_w, cell_h, title, body or "—")

    # --- 4. フッター（右下：Value Create） ---
    footer_box = slide.shapes.add_textbox(
        SLIDE_W - Inches(3.0), SLIDE_H - Inches(0.45), Inches(2.7), Inches(0.3)
    )
    _set_textbox_text(
        footer_box,
        "Value Create",
        font_size=10,
        bold=True,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def _draw_cell(slide, x, y, w, h, title: str, body: str) -> None:
    """1つのセクションセルを描画する。本文は箇条書き + **bold** インライン対応。
    ダーク基調: surface.raised 背景 + border.muted 境界 + 白テキスト + イエロー左アクセント。"""
    # 背景（raised-gray の角丸矩形 + 淡い境界線）
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.adjustments[0] = 0.06  # radius 13px 相当に近いラウンド
    bg.fill.solid()
    bg.fill.fore_color.rgb = SURFACE_RAISED
    bg.line.color.rgb = BORDER_MUTED
    bg.line.width = Pt(0.5)

    # 左端に黄色のバーチカルアクセント
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_YELLOW

    # タイトル（白）
    title_tb = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.1), w - Inches(0.3), Inches(0.4)
    )
    _set_textbox_text(
        title_tb, title, font_size=13, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT
    )

    # 本文（箇条書き + インライン太字、白）。4 bullet まで収めるため 11pt に抑える。
    body_tb = slide.shapes.add_textbox(
        x + Inches(0.2),
        y + Inches(0.52),
        w - Inches(0.3),
        h - Inches(0.62),
    )
    _write_body_bullets(body_tb, body, font_size=11, color=TEXT_PRIMARY)


_BOLD_SPLIT = re.compile(r"(\*\*[^*]+\*\*)")


def _parse_bold_segments(line: str) -> list[tuple[str, bool]]:
    """'- **Keyword**: desc' → [('Keyword', True), (': desc', False)] のように分解する。
    先頭の箇条書きプレフィックス（`- ` / `* `）は呼び出し側で除去しておくこと。"""
    segments: list[tuple[str, bool]] = []
    for part in _BOLD_SPLIT.split(line):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**") and len(part) >= 4:
            segments.append((part[2:-2], True))
        else:
            segments.append((part, False))
    return segments


def _write_body_bullets(textbox, body: str, *, font_size: int, color: RGBColor) -> None:
    """本文テキストを箇条書き + インライン太字で描画する。
    - 行頭 `- ` または `* ` は視覚的な `• ` に置換
    - `**text**` は太字 run として分割
    - 空行はスキップ"""
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.clear()

    lines = [ln for ln in (body or "").split("\n") if ln.strip()]
    if not lines:
        lines = ["—"]

    for i, raw in enumerate(lines):
        stripped = raw.strip()
        # 箇条書きプレフィックスを視覚的な `• ` に置換
        if stripped.startswith("- "):
            stripped = "• " + stripped[2:]
        elif stripped.startswith("* "):
            stripped = "• " + stripped[2:]

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)

        for seg_text, is_bold in _parse_bold_segments(stripped):
            run = p.add_run()
            run.text = seg_text
            run.font.size = Pt(font_size)
            run.font.bold = is_bold
            run.font.color.rgb = color
            run.font.name = FONT_FAMILY


def _set_shape_text(
    shape,
    text: str,
    *,
    font_size: int,
    bold: bool,
    color: RGBColor,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    left_margin=None,
) -> None:
    """図形内のテキストを設定する。"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if left_margin is not None:
        tf.margin_left = left_margin
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    _write_paragraph(tf, text, font_size=font_size, bold=bold, color=color, align=align)


def _set_textbox_text(
    textbox,
    text: str,
    *,
    font_size: int,
    bold: bool,
    color: RGBColor,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    word_wrap: bool = True,
) -> None:
    tf = textbox.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    _write_paragraph(tf, text, font_size=font_size, bold=bold, color=color, align=align)


def _write_paragraph(tf, text: str, *, font_size: int, bold: bool, color: RGBColor, align) -> None:
    """text_frame 内の paragraph 群を text で置き換える。
    改行を \n で判定し複数段落を生成する。Noto Sans Japanese を指定。"""
    # 既存の段落をクリア（最初の段落は残す）
    tf.clear()
    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT_FAMILY
